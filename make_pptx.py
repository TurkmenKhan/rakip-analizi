"""
ISP Intel - CGO Sunumu
TurkNet · Kiraz Çiçeği × Derin Mor Teması
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.abspath(__file__))

C_DARK    = RGBColor(0x12, 0x07, 0x1E)
C_CARD    = RGBColor(0x1E, 0x0D, 0x35)
C_ACCENT  = RGBColor(0xFF, 0xB7, 0xC5)
C_ACCENT2 = RGBColor(0xD4, 0x5E, 0x8B)
C_LAVNDR  = RGBColor(0xB5, 0x7E, 0xDC)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY   = RGBColor(0xE8, 0xD5, 0xEC)
C_MGRAY   = RGBColor(0x8B, 0x7A, 0x9B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    slide = prs.slides.add_slide(BLANK)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C_DARK
    return slide

def box(slide, x, y, w, h, color=None):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.line.fill.background()
    if color:
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
    else:
        sh.fill.background()
    return sh

def txt(slide, text, x, y, w, h, size=24, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return txb

def line(slide, x, y, w, color=C_ACCENT, h=0.05):
    box(slide, x, y, w, h, color)

def kpi(slide, x, y, w, h, label, value, sub="", color=C_ACCENT):
    box(slide, x, y, w, h, C_CARD)
    line(slide, x, y, w, color, 0.05)
    txt(slide, label, x+0.2, y+0.18, w-0.4, 0.45, size=11, color=C_MGRAY)
    txt(slide, value, x+0.2, y+0.55, w-0.4, 1.0,  size=38, bold=True, color=color)
    if sub:
        txt(slide, sub, x+0.2, y+1.45, w-0.4, 0.4, size=13, color=C_MGRAY)

def bullet3(slide, x, y, w, items, color=C_ACCENT):
    """Maksimum 3 madde, büyük font."""
    for i, item in enumerate(items[:3]):
        cy = y + i * 0.9
        box(slide, x, cy+0.15, 0.06, 0.5, color)
        txt(slide, item, x+0.25, cy, w-0.3, 0.75, size=20, color=C_LGRAY)

def add_cherry_petals(slide):
    positions = [
        (0.12, 0.15, 0.28, 0.20), (12.30, 0.25, 0.32, 0.22),
        (0.35, 6.65, 0.22, 0.16), (12.00, 6.35, 0.30, 0.20),
        (6.25, 0.08, 0.25, 0.17), (10.75, 3.35, 0.20, 0.14),
        (1.35, 3.85, 0.22, 0.16), (5.80, 7.10, 0.18, 0.13),
        (9.50, 0.05, 0.20, 0.14),
    ]
    for (x, y, w, h) in positions:
        sh = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = C_ACCENT
        sh.line.fill.background()
    for (x, y, w, h) in positions[::2]:
        sh = slide.shapes.add_shape(9,
            Inches(x+0.05), Inches(y+0.06), Inches(w*0.65), Inches(h*0.65))
        sh.fill.solid()
        sh.fill.fore_color.rgb = C_ACCENT2
        sh.line.fill.background()

def turknet_logo(slide, x, y, w=2.6):
    logo_path = os.path.join(BASE, "turknet_logo.png")
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(x), Inches(y), Inches(w))
    else:
        txt(slide, "turknet", x, y, w, 0.6, size=28, color=C_WHITE, align=PP_ALIGN.CENTER)

def slide_label(slide, text, color=C_ACCENT):
    txt(slide, text, 0.5, 0.3, 12.0, 0.45, size=11, bold=True, color=color)

def slide_title(slide, text):
    txt(slide, text, 0.5, 0.75, 12.3, 1.4, size=40, bold=True, color=C_WHITE)

def divider(slide, y=2.05):
    line(slide, 0.5, y, 12.3, C_ACCENT, 0.04)


# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 1 — KAPAK
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_cherry_petals(s)

# Sol koyu panel
box(s, 0, 0, 5.8, 7.5, C_CARD)
box(s, 0, 0, 0.08, 7.5, C_ACCENT)

turknet_logo(s, 0.55, 0.9, 2.8)
line(s, 0.55, 1.75, 4.8, C_ACCENT)

txt(s, "ISP INTEL",
    0.55, 2.0, 5.0, 1.2, size=52, bold=True, color=C_WHITE)
txt(s, "Rakip İzleme & Pazar Analizi",
    0.55, 3.2, 5.0, 0.7, size=20, color=C_ACCENT)
txt(s, "Türkiye İnternet Sektörü · Haziran 2026",
    0.55, 4.0, 5.0, 0.5, size=14, color=C_MGRAY)

# Sağ panel — 4 büyük KPI
txt(s, "CANLI SİSTEM", 6.2, 0.9, 6.8, 0.5, size=12, bold=True, color=C_ACCENT)
line(s, 6.2, 1.45, 6.8, C_ACCENT)

kpi(s, 6.2,  1.65, 3.2, 2.0, "İZLENEN ISS",     "22",   "aktif sağlayıcı",   C_ACCENT)
kpi(s, 9.6,  1.65, 3.5, 2.0, "AKTİF PAKET",      "230",  "güncel fiyat",      C_LAVNDR)
kpi(s, 6.2,  3.85, 3.2, 2.0, "GÜNCELLEME",       "30dk", "otomatik döngü",    C_ACCENT2)
kpi(s, 9.6,  3.85, 3.5, 2.0, "TAKİP NOKTASI",    "193",  "izlenen URL",       C_ACCENT)

txt(s, "GİZLİ — YÖNETİM SUNUMU", 6.2, 7.1, 6.8, 0.35,
    size=9, color=C_MGRAY, italic=True, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 2 — PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_cherry_petals(s)
box(s, 0, 0, 0.08, 7.5, C_ACCENT2)

slide_label(s, "PROBLEM", C_ACCENT2)
slide_title(s, "Rakipler ne yapıyor?\nBilmiyorduk.")
divider(s, 2.15)

bullet3(s, 0.5, 2.45, 7.5, [
    "128 ISS'in fiyatlarını günlük takip etmek insan gücüyle imkânsız.",
    "Rakip bir indirim veya zam yaptığında günlerce haberdar olamıyorduk.",
    "Fiyat stratejisi veriye değil, sezgiye dayanıyordu.",
], color=C_ACCENT2)

# Sağ — büyük rakam
box(s, 8.5, 2.0, 4.5, 5.0, C_CARD)
line(s, 8.5, 2.0, 4.5, C_ACCENT2)
txt(s, "128", 8.5, 2.4, 4.5, 2.0, size=96, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
txt(s, "rakip ISS", 8.5, 4.3, 4.5, 0.6, size=20, color=C_MGRAY, align=PP_ALIGN.CENTER)
line(s, 8.8, 5.1, 3.9, C_CARD)
txt(s, "takip edilmesi gereken", 8.5, 5.3, 4.5, 0.5, size=14, color=C_MGRAY, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 3 — ÇÖZÜM
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_cherry_petals(s)
box(s, 0, 0, 0.08, 7.5, C_LAVNDR)

slide_label(s, "ÇÖZÜM", C_LAVNDR)
slide_title(s, "Tam otomatik.\n7/24 kesintisiz.")
divider(s, 2.15)

# 4 adım — büyük kutular (oklar shape ile)
steps = [
    ("Web İzleme",      "193 URL\nOtomatik Taranır",      C_ACCENT),
    ("Değişim Tespiti", "Sadece Fark\nAnaliz Edilir",      C_ACCENT2),
    ("Yapay Zeka",      "Paket Bilgisi\nÇıkarılır",        C_LAVNDR),
    ("Dashboard",       "Anlık Görüntü\n& Bildirim",       C_ACCENT),
]
BOX_W = 2.85
GAP   = 0.45
for i, (title, desc, col) in enumerate(steps):
    cx = 0.5 + i * (BOX_W + GAP)
    box(s, cx, 2.45, BOX_W, 3.2, col)
    txt(s, str(i+1), cx+0.2, 2.58, BOX_W-0.3, 0.75, size=28, bold=True,  color=C_DARK)
    txt(s, title,    cx+0.2, 3.32, BOX_W-0.3, 0.65, size=15, bold=True,  color=C_DARK)
    txt(s, desc,     cx+0.2, 3.95, BOX_W-0.3, 1.30, size=18, bold=False, color=C_DARK)
    if i < 3:
        # İnce yatay bağlantı çizgisi + üçgen ok (shape)
        mid_y = 2.45 + 3.2/2 - 0.03
        box(s, cx+BOX_W, mid_y, GAP*0.55, 0.06, C_MGRAY)
        # Üçgen ok: right-arrow shape (type 13)
        from pptx.util import Emu
        arr = s.shapes.add_shape(13,
            Inches(cx+BOX_W+GAP*0.5), Inches(mid_y-0.12),
            Inches(GAP*0.45), Inches(0.30))
        arr.fill.solid()
        arr.fill.fore_color.rgb = C_MGRAY
        arr.line.fill.background()

txt(s, "İnsan müdahalesi gerektirmez. Sıfır manuel iş.",
    0.5, 6.0, 12.3, 0.6, size=18, color=C_LAVNDR, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 4 — PAZAR GÖRÜNÜMÜ (Dashboard screenshots)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_cherry_petals(s)
box(s, 0, 0, 0.08, 7.5, C_ACCENT)

slide_label(s, "PLATFORM", C_ACCENT)
slide_title(s, "Tek ekranda tüm pazar.")
divider(s, 2.15)

img_genel = os.path.join(BASE, "stitch", "stitch", "genel_bak_isp_intel", "screen.png")
img_karsi = os.path.join(BASE, "stitch", "stitch", "kar_la_t_rma_isp_intel", "screen.png")

if os.path.exists(img_genel) and os.path.exists(img_karsi):
    s.shapes.add_picture(img_genel, Inches(0.5),  Inches(2.3), Inches(5.9))
    s.shapes.add_picture(img_karsi, Inches(6.9),  Inches(2.3), Inches(5.9))
    txt(s, "Genel Bakış", 0.5, 6.55, 5.9, 0.45, size=14, color=C_MGRAY, align=PP_ALIGN.CENTER)
    txt(s, "Karşılaştırma", 6.9, 6.55, 5.9, 0.45, size=14, color=C_MGRAY, align=PP_ALIGN.CENTER)
else:
    txt(s, "[ Dashboard Ekran Görüntüsü ]",
        0.5, 3.5, 12.3, 1.0, size=24, color=C_MGRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 5 — KİLİT BULGULAR
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_cherry_petals(s)
box(s, 0, 0, 0.08, 7.5, C_ACCENT2)

slide_label(s, "KİLİT BULGULAR", C_ACCENT2)
slide_title(s, "Piyasa ne söylüyor?")
divider(s, 2.15)

# 3 büyük insight kutusu
insights = [
    ("100 Mbps\nPazarı",
     "En yoğun rekabet bölgesi.\nMin 345 TL — Max 1.335 TL\narasında geniş fark.",
     C_ACCENT),
    ("1 Gbps\nFırsatı",
     "Piyasada 500 TL'den başlıyor.\nYüksek büyüme segmenti,\nrakipler hızlı giriyor.",
     C_LAVNDR),
    ("Kampanya\nHızı",
     "Rakipler ortalama ayda\n2-3 kampanya değişikliği yapıyor.\nTepki süresi kritik.",
     C_ACCENT2),
]
for i, (title, body, col) in enumerate(insights):
    cx = 0.5 + i * 4.2
    box(s, cx, 2.3, 3.9, 4.7, C_CARD)
    line(s, cx, 2.3, 3.9, col)
    txt(s, title, cx+0.25, 2.5,  3.5, 1.1, size=22, bold=True, color=col)
    txt(s, body,  cx+0.25, 3.65, 3.5, 3.1, size=16, color=C_LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 6 — İŞ DEĞERİ
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_cherry_petals(s)
box(s, 0, 0, 0.08, 7.5, C_LAVNDR)

slide_label(s, "İŞ DEĞERİ", C_LAVNDR)
slide_title(s, "Ne kazandık?")
divider(s, 2.15)

# 2 büyük kazanım + 1 fırsat
gains = [
    ("Sıfır\nManuel İş",
     "Haftada saatler harcanan manuel takip tamamen ortadan kalktı.",
     C_LAVNDR),
    ("30 Dakika\nTepki Süresi",
     "Eskiden günler alıyordu. Şimdi rakip hareket ettiği anda bildirim alıyoruz.",
     C_ACCENT),
    ("Veri Bazlı\nStrateji",
     "Fiyat kararları artık sezgiye değil, 230 paketin gerçek zamanlı analizine dayanıyor.",
     C_ACCENT2),
]
for i, (title, body, col) in enumerate(gains):
    cx = 0.5 + i * 4.2
    box(s, cx, 2.3, 3.9, 4.7, C_CARD)
    line(s, cx, 2.3, 3.9, col)
    txt(s, title, cx+0.25, 2.5,  3.5, 1.2, size=26, bold=True, color=col)
    txt(s, body,  cx+0.25, 3.8,  3.5, 3.0, size=16, color=C_LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 7 — SONUÇ & SONRAKİ ADIMLAR
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_cherry_petals(s)
box(s, 0, 0, 0.08, 7.5, C_ACCENT)

slide_label(s, "SONUÇ", C_ACCENT)
slide_title(s, "Sistem canlıda.\nSıradaki adım ne?")
divider(s, 2.15)

# Sol — tamamlananlar (3 madde)
txt(s, "TAMAMLANDI", 0.5, 2.35, 5.8, 0.45, size=12, bold=True, color=C_LAVNDR)
bullet3(s, 0.5, 2.9, 5.8, [
    "22 ISS, 193 URL — 7/24 otomatik izlemede.",
    "230 aktif paket anlık takipte, fiyat değişim bildirimi aktif.",
    "Web dashboard — kurulum gerektirmez, her cihazdan erişim.",
], color=C_LAVNDR)

# Sağ — sonraki adımlar (3 madde)
box(s, 7.0, 2.2, 6.0, 4.8, C_CARD)
line(s, 7.0, 2.2, 6.0, C_ACCENT2)
txt(s, "SONRAKİ ADIMLAR", 7.2, 2.38, 5.6, 0.45, size=12, bold=True, color=C_ACCENT2)
bullet3(s, 7.2, 2.95, 5.5, [
    "E-posta & Telegram bildirim entegrasyonu.",
    "Haftalık otomatik yönetim raporu.",
    "Fiyat trend analizi ve tahmin modeli.",
], color=C_ACCENT2)

# Alt şerit
box(s, 0.5, 6.55, 12.3, 0.75, C_CARD)
line(s, 0.5, 6.55, 12.3, C_ACCENT, 0.04)
turknet_logo(s, 0.75, 6.63, 2.0)
txt(s, "Sistem şu an aktif — gerçek zamanlı veri için dashboard'u ziyaret edin.",
    3.0, 6.68, 9.5, 0.5, size=14, color=C_LGRAY)


# ── Kaydet ─────────────────────────────────────────────────────────────────────
out = r"C:\Users\erhan\OneDrive\Masaüstü\ISP_Intel_CGO_Sunumu.pptx"
prs.save(out)
print(f"Kaydedildi: {out}")
print(f"Slayt sayısı: {len(prs.slides)}")
